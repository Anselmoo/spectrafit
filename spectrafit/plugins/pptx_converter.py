"""Convert the lock file to a powerpoint presentation."""
import argparse

from pathlib import Path
from typing import Any
from typing import Dict
from typing import MutableMapping
from typing import Type

import pandas as pd
import tomli

from pptx import Presentation
from pptx.util import Pt
from spectrafit.api.pptx_model import PPTXDataAPI
from spectrafit.api.pptx_model import PPTXLayoutAPI
from spectrafit.api.pptx_model import PPTXPositionAPI
from spectrafit.api.pptx_model import PPTXRatioAPI
from spectrafit.api.pptx_model import PPTXStructureAPI
from spectrafit.plugins.converter import Converter


class PPTXElements:
    """Generate a powerpoint presentation from a the spectrafit output."""

    def __init__(self, slide: Type[Any]) -> None:
        """Create a powerpoint presentation from a the spectrafit output.

        Args:
            slide (type): The slide of the powerpoint presentation.
        """
        self.slide = slide

    def create_textbox(
        self,
        text: str,
        position: PPTXPositionAPI,
        font_size: Pt = Pt(16),
    ) -> None:
        """Create a textbox from the input file.

        Args:
            text (str): The text of the textbox.
            position (PPTXPositionAPI): The position of the textbox in the powerpoint
                presentation.
            font_size (Pt): The font size of the textbox. Defaults to Pt(16).
        """
        tx_box = self.slide.shapes.add_textbox(
            left=position.left,
            top=position.top,
            width=position.width,
            height=position.height,
        )

        tf = tx_box.text_frame
        tf.text = text
        tf.size = font_size

    def create_figure(
        self,
        fname: Path,
        position_figure: PPTXPositionAPI,
        position_textbox: PPTXPositionAPI,
        text: str,
        font_size: Pt,
    ) -> None:
        """Create a figure from the input file.

        Args:
            fname (Path): The temporay filename of the figure.
            position_figure (PPTXPositionAPI): The position of the figure in the
                powerpoint presentation.
            position_textbox (PPTXPositionAPI): The position of the textbox in the
                powerpoint presentation.
            text (str): The text of the textbox.
            font_size (Pt): The font size of the textbox.
        """
        self.slide.shapes.add_picture(
            str(fname),
            left=position_figure.left,
            top=position_figure.top,
            width=position_figure.width,
            height=position_figure.height,
        )
        self.create_textbox(
            text=text,
            position=position_textbox,
            font_size=font_size,
        )

    def create_table(
        self,
        df: pd.DataFrame,
        position_table: PPTXPositionAPI,
        transpose: bool,
        index_hidden: bool,
        text: str,
        position_textbox: PPTXPositionAPI,
        font_size: Pt = Pt(12),
    ) -> None:
        """Create a table from the input file.

        Args:
            df (pd.DataFrame): The data of the table.
                powerpoint presentation.
            position_table (PPTXPositionAPI): The position of the table in the
                powerpoint presentation.
            transpose (bool): Transpose the table for row-like presentation in
                powerpoint presentation.
            index_hidden (bool): Hide the index of the table in the powerpoint
                presentation.
            text (str): The text of the table legend.
            position_textbox (PPTXPositionAPI): The position of the textbox in the
                powerpoint presentation.
            font_size (Pt, optional): The font size of the table. Defaults to Pt(12).
        """
        df = df.round(2)
        if transpose:
            df = df.transpose()
        self.extract_table(
            df=df, position_table=position_table, index_hidden=index_hidden
        )

        self.create_textbox(
            text=text,
            position=PPTXPositionAPI(
                left=position_textbox.left,
                top=position_textbox.top,
                width=position_textbox.width,
                height=position_textbox.height,
            ),
            font_size=font_size,
        )

    def extract_table(
        self,
        df: pd.DataFrame,
        position_table: PPTXPositionAPI,
        index_hidden: bool,
    ) -> None:
        """Create a table from the input file.

        Args:
            df (pd.DataFrame): The data of the table(s) in the powerpoint presentation.
            position_table (PPTXPositionAPI): The position of the table in the
                powerpoint presentation.
            index_hidden (bool): Hide the index of the table in the powerpoint
                presentation.
        """
        rows, cols = df.shape
        table = self.slide.shapes.add_table(
            rows=rows + 1,
            cols=cols + (not index_hidden),
            left=position_table.left,
            top=position_table.top,
            width=position_table.width,
            height=position_table.height,
        )
        table.table.cell(0, 0).text = ""
        if index_hidden:
            for i, col in enumerate(df.columns):
                table.table.cell(0, i).text = str(col)
                for j, value in enumerate(df[col]):
                    table.table.cell(j + 1, i).text = str(value)
        else:
            for i, index in enumerate(df.index, start=1):
                table.table.cell(i, 0).text = str(index)
            for i, col in enumerate(df.columns, start=1):
                table.table.cell(0, i).text = col
                for j, value in enumerate(df[col]):
                    table.table.cell(j + 1, i).text = str(value)

    def create_credit(
        self,
        fname: Path,
        text: str,
        position_logo: PPTXPositionAPI,
        position_text: PPTXPositionAPI,
        font_size: Pt = Pt(14),
    ) -> None:
        """Create a credit for spectrafit.

        Args:
            fname (Path): The temporay filename of the figure.
            text (str): The text of the credit.
            position_logo (PPTXPositionAPI): The position of the logo in the powerpoint
                presentation.
            position_text (PPTXPositionAPI): The position of the text in the powerpoint
                presentation.
            font_size (Pt): The font size of the textbox. Defaults to Pt(14).
        """
        self.create_figure(
            fname=fname,
            position_figure=position_logo,
            text=text,
            position_textbox=position_text,
            font_size=font_size,
        )

    def create_title(self, text: str, position: PPTXPositionAPI) -> None:
        """Create a title from the input file.

        Args:
            text (str): The text of the title.
            position (PPTXPositionAPI): The position of the title in the powerpoint
                presentation.
        """
        title = self.slide.shapes.title
        title.text = text
        title.left = position.left
        title.top = position.top
        title.width = position.width
        title.height = position.height

    def create_subtitle(self, text: str, position: PPTXPositionAPI, index: int) -> None:
        """Create a subtitle from the input file.

        Args:
            text (str): The text of the subtitle.
            position (PPTXPositionAPI): The position of the subtitle in the powerpoint
                presentation.
            index (int): The index of the subtitle in the powerpoint presentation.
        """
        subtitle = self.slide.placeholders[index]
        subtitle.text = text
        subtitle.left = position.left
        subtitle.top = position.top
        subtitle.width = position.width
        subtitle.height = position.height


class PPTXLayout(PPTXElements):
    """Generate a powerpoint presentation from a the spectrafit output."""

    def __init__(
        self, ratio: PPTXRatioAPI, structure: PPTXStructureAPI, fname: Path
    ) -> None:
        """Create a powerpoint presentation from a the spectrafit output.

        Args:
            ratio (PPTXRatioAPI): The ratio of the powerpoint presentation.
            structure (PPTXStructureAPI): The structure of the powerpoint presentation.
            fname (Path): The temporay filename of the powerpoint presentation.
        """
        self.ratio = ratio
        self.structure = structure
        self.fname = fname

        self._initialiaze()
        super().__init__(self.slide)

    def _initialiaze(self) -> None:
        """Initialize the powerpoint presentation."""
        self.prs = Presentation()
        self.prs.slide_width = self.ratio.width
        self.prs.slide_height = self.ratio.height
        self.slide = self.prs.slides.add_slide(self.prs.slide_layouts[3])

    def top_element(self) -> None:
        """Create the top element of the powerpoint presentation."""
        self.create_title(
            text=self.structure.header.text, position=self.structure.header.position
        )

    def lefr_element(self) -> None:
        """Create the left element of the powerpoint presentation."""
        self.create_subtitle(
            text=self.structure.sub_title_left.text,
            position=self.structure.sub_title_left.position,
            index=self.structure.sub_title_left.index,
        )
        self.create_figure(
            fname=self.structure.sub_title_left.figure.fname,
            position_figure=self.structure.sub_title_left.figure.position,
            text=self.structure.sub_title_left.figure.description.text,
            position_textbox=self.structure.sub_title_left.figure.description.position,
            font_size=self.structure.sub_title_left.figure.description.font_size,
        )

    def right_element(self) -> None:
        """Create the right element of the powerpoint presentation."""
        self.create_subtitle(
            text=self.structure.sub_title_right.text,
            position=self.structure.sub_title_right.position,
            index=self.structure.sub_title_right.index,
        )
        self.create_table(
            df=self.structure.sub_title_right.table_1.df,
            position_table=self.structure.sub_title_right.table_1.position,
            transpose=self.structure.sub_title_right.table_1.transpose,
            index_hidden=self.structure.sub_title_right.table_1.index_hidden,
            text=self.structure.sub_title_right.table_1.description.text,
            position_textbox=(
                self.structure.sub_title_right.table_1.description.position
            ),
            font_size=self.structure.sub_title_right.table_1.description.font_size,
        )
        self.create_table(
            df=self.structure.sub_title_right.table_2.df,
            position_table=self.structure.sub_title_right.table_2.position,
            transpose=self.structure.sub_title_right.table_2.transpose,
            index_hidden=self.structure.sub_title_right.table_2.index_hidden,
            text=self.structure.sub_title_right.table_2.description.text,
            position_textbox=(
                self.structure.sub_title_right.table_2.description.position
            ),
            font_size=self.structure.sub_title_right.table_2.description.font_size,
        )
        self.create_table(
            df=self.structure.sub_title_right.table_3.df,
            position_table=self.structure.sub_title_right.table_3.position,
            transpose=self.structure.sub_title_right.table_3.transpose,
            index_hidden=self.structure.sub_title_right.table_3.index_hidden,
            text=self.structure.sub_title_right.table_3.description.text,
            position_textbox=(
                self.structure.sub_title_right.table_3.description.position
            ),
            font_size=self.structure.sub_title_right.table_3.description.font_size,
        )
        self.create_credit(
            fname=self.structure.sub_title_right.credit.fname,
            text=self.structure.sub_title_right.credit.description.text,
            position_logo=self.structure.sub_title_right.credit.position,
            position_text=self.structure.sub_title_right.credit.description.position,
            font_size=self.structure.sub_title_right.credit.description.font_size,
        )

    def save(self) -> None:
        """Save the powerpoint presentation."""
        self.prs.save(str(self.fname))

    def __call__(self) -> None:
        """Create the powerpoint presentation."""
        self.top_element()
        self.lefr_element()
        self.right_element()
        self.save()


class PPTXConverter(Converter):
    """Generate a powerpoint presentation from a the spectrafit output.

    Attributes:
        pixel_size (Dict[str, Dict[str, int]]): The pixel size of the powerpoint
            presentation.
    """

    pixel_size = PPTXLayoutAPI.pptx_formats.keys()

    def get_args(self) -> Dict[str, Any]:
        """Get the arguments from the command line.

        Returns:
            Dict[str, Any]: Return the input file arguments as a dictionary without
                additional information beyond the command line arguments.
        """
        parse = argparse.ArgumentParser(
            description="Converter for 'SpectraFit' from *.lock output to a "
            "PowerPoint presentation.",
            usage="%(prog)s [options] infile",
        )
        parse.add_argument(
            "infile",
            type=Path,
            help="Filename of the *.lock file to convert to a powerpoint presentation.",
        )
        parse.add_argument(
            "-f",
            "--file-format",
            help="File format of the PowerPoint presentation. Default is '16:9'.",
            type=str,
            default="16:9",
            choices=self.pixel_size,
        )
        return vars(parse.parse_args())

    @staticmethod
    def convert(infile: Path, file_format: str) -> MutableMapping[str, Any]:
        """Convert the lock file to a powerpoint presentation.

        Args:
            infile (Path): The input file of the as a path object.
            file_format (str): The powerpoint presentation file format.

        Raises:
            ValueError: If the file format is not supported.
            ValueError: If the input file is not a `*.lock` file.

        Returns:
            MutableMapping[str, Any]: The converted file as a dictionary.
        """
        if file_format not in PPTXConverter.pixel_size:
            raise ValueError(
                f"File format '{file_format}' is not supported;"
                f"it must be one of {PPTXConverter.pixel_size}"
            )

        if infile.suffix != ".lock":
            raise ValueError(
                f"File format '{infile.suffix}' is not supported; it must be '.lock'"
            )

        with open(infile, "rb") as f:
            data = PPTXDataAPI(**tomli.load(f))

        return {file_format: data}

    def save(self, data: Any, fname: Path, export_format: str) -> None:
        """Save the powerpoint presentation."""
        pptx_layout = PPTXLayoutAPI(
            export_format, data=data[export_format]
        ).get_pptx_layout()

        PPTXLayout(
            ratio=pptx_layout.ratio,
            structure=pptx_layout.structure,
            fname=Path(f"{fname.stem}_{export_format.replace(':', '_')}.pptx"),
        )()

    def __call__(self) -> None:
        """Convert the lock file to a powerpoint presentation."""
        args = self.get_args()
        data = self.convert(args["infile"], args["file_format"])
        self.save(data, args["infile"], args["file_format"])


def command_line_runner() -> None:
    """Command line interface for the converter plugin."""
    PPTXConverter()()
